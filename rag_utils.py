# import os
# import tempfile

# from langchain_community.document_loaders import UnstructuredFileLoader
# from langchain_text_splitters import RecursiveCharacterTextSplitter
# from langchain_huggingface import HuggingFaceEmbeddings
# from langchain_community.vectorstores import FAISS

# from langchain_core.prompts import ChatPromptTemplate
# from langchain_groq import ChatGroq
# from langchain_classic.chains.combine_documents import create_stuff_documents_chain
# from langchain_classic.chains.retrieval import create_retrieval_chain


# # File handling + extraction

# def save_temp(uploaded):
#     ext = os.path.splitext(uploaded.name)[1]
#     tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
#     tmp.write(uploaded.getbuffer())
#     return tmp.name


# def extract_chunks(file_path, chunk_size=1200, chunk_overlap=200):
#     loader = UnstructuredFileLoader(file_path)
#     docs = loader.load()

#     splitter = RecursiveCharacterTextSplitter(
#         chunk_size=chunk_size,
#         chunk_overlap=chunk_overlap
#     )

#     return splitter.split_documents(docs)



# # Vector DB (Embeddings + FAISS)

# def build_vector_db(all_chunks):
#     embedder = HuggingFaceEmbeddings(
#         model_name="sentence-transformers/all-MiniLM-L6-v2"
#     )
#     return FAISS.from_documents(all_chunks, embedder)



# # RAG Chain Builder

# def build_rag_chain(retriever, api_key):

#     prompt = ChatPromptTemplate.from_template("""
# Use ONLY the context below.

# Context:
# {context}

# Instruction:
# {input}

# Response:
# """)

#     llm = ChatGroq(
#         api_key=api_key,
#         model="llama-3.1-8b-instant",
#         temperature=0.5
#     )

#     combine = create_stuff_documents_chain(llm=llm, prompt=prompt)

#     return create_retrieval_chain(
#         retriever=retriever,
#         combine_docs_chain=combine



#     )# rag_utils.py


import os
import time
import tempfile
import mimetypes
import json
import logging
from concurrent.futures import ThreadPoolExecutor, as_completed
from PIL import Image
import pytesseract
from aud_vid_utils import transcribe_fast

# Fast extractors
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

# textract is optional (good fallback for many formats)
try:
    import textract
except Exception:
    textract = None

# LangChain / embeddings / vectorstore
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.prompts import ChatPromptTemplate
from langchain_groq import ChatGroq
from langchain_classic.chains.combine_documents import create_stuff_documents_chain
from langchain_classic.chains.retrieval import create_retrieval_chain

# Last-resort fallback loader (only used if everything else fails)
try:
    from langchain_community.document_loaders import UnstructuredFileLoader
except Exception:
    UnstructuredFileLoader = None

# Optional: Document wrapper for FAISS
try:
    from langchain_core.documents import Document as LCDocument
except Exception:
    # minimal fallback for old/new versions -- FAISS.from_documents expects objects with page_content
    class LCDocument:
        def __init__(self, page_content, metadata=None):
            self.page_content = page_content
            self.metadata = metadata or {}

# Logging
logger = logging.getLogger("rag_utils")
logger.setLevel(logging.INFO)


# -------------------------
# TEMP FILE HANDLER
# -------------------------
def save_temp(uploaded):
    """Save an uploaded Streamlit file to a local temp file and return its path."""
    ext = os.path.splitext(uploaded.name)[1]
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
    tmp.write(uploaded.getbuffer())
    tmp.flush()
    tmp.close()
    return tmp.name


# -------------------------
# FILE TYPE DETECTION
# -------------------------
def detect_file_type(path):
    """Return (ext, mime) for a path. Ext lowercased (including leading dot)."""
    ext = os.path.splitext(path)[1].lower()
    mime, _ = mimetypes.guess_type(path)
    return ext, mime


# -------------------------
# FAST EXTRACTORS
# Each extractor returns a single big text string for that file.
# -------------------------
def extract_pdf(path):
    if fitz is None:
        raise RuntimeError("PyMuPDF (fitz) is required for PDF extraction but not installed.")
    doc = fitz.open(path)
    parts = []
    for page in doc:
        parts.append(page.get_text())
    return "\n".join(parts)


def extract_docx(path):
    if Document is None:
        raise RuntimeError("python-docx is required for DOCX extraction but not installed.")
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])


def extract_pptx(path):
    if Presentation is None:
        raise RuntimeError("python-pptx is required for PPTX extraction but not installed.")
    prs = Presentation(path)
    slides = []
    for sli, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        slides.append("\n".join(texts))
    return "\n\n--- SLIDE BREAK ---\n\n".join(slides)


def extract_image(path):
    """
    Fast image → text using pytesseract (2–3 sec).
    Auto downscales huge images for speed.
    """
    img = Image.open(path)

    # downscale if extremely large ( > 5K x 5K )
    max_dim = 3500
    if img.width > max_dim or img.height > max_dim:
        scale = max_dim / max(img.width, img.height)
        new_size = (int(img.width * scale), int(img.height * scale))
        img = img.resize(new_size)

    text = pytesseract.image_to_string(img)
    return text


def extract_text(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_textract(path):
    if textract is None:
        return None
    try:
        raw = textract.process(path)
        return raw.decode("utf-8", errors="ignore")
    except Exception:
        return None


def extract_unstructured(path):
    if UnstructuredFileLoader is None:
        raise RuntimeError("UnstructuredFileLoader is not available as a fallback.")
    docs = UnstructuredFileLoader(path).load()
    return "\n".join([d.page_content for d in docs])


# -------------------------
# MASTER EXTRACTOR (fast + universal)
# -------------------------
def fast_extract(path):
    """
    Fast extractor that:
    1) tries specialized extractors for common formats
    2) falls back to textract for many other types
    3) falls back to Unstructured as a last resort
    Returns: extracted text string
    """
    ext, mime = detect_file_type(path)

    # Fast specialized loaders
    try:
        if ext == ".pdf":
            return extract_pdf(path)
        if ext == ".docx":
            return extract_docx(path)
        if ext == ".pptx":
            return extract_pptx(path)
        if ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff"):
            return extract_image(path)        
        if ext in (".mp3", ".wav", ".m4a", ".flac", ".mp4", ".mov", ".avi"):
            return transcribe_fast(path)
        if ext in (".txt", ".md", ".csv", ".json"):
            return extract_text(path)
        
    except Exception as e:
        logger.warning(f"Specialized extractor for {ext} failed: {e}")

    # universal fallback: textract
    raw = extract_textract(path)
    if raw:
        return raw

    # final fallback: Unstructured
    try:
        return extract_unstructured(path)
    except Exception as e:
        logger.error(f"All extractors failed for {path}: {e}")
        return ""


# FAST CHUNKING

def fast_chunk(text, chunk_size=1200, chunk_overlap=200):
    """
    Return a list of text chunks (strings).
    Uses LangChain's RecursiveCharacterTextSplitter but returns plain strings for speed.
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap
    )
    # split_text is faster and returns list[str]
    return splitter.split_text(text)



# EXTRACT + CHUNK (per file)

def extract_chunks(path, chunk_size=1200, chunk_overlap=200):
    """
    Extract & chunk a single file path. Returns list[str].
    """
    text = fast_extract(path)
    if not text:
        return []
    return fast_chunk(text, chunk_size=chunk_size, chunk_overlap=chunk_overlap)



# MULTI-FILE PARALLEL PROCESSING + PROGRESS + BENCHMARK LOGS

def process_files(
    file_paths,
    max_workers=4,
    chunk_size=1200,
    chunk_overlap=200,
    progress_callback=None,
):
    """
    Process a list of local file paths in parallel:
      - extract
      - chunk

    progress_callback is optional and called with a dict:
      {
        "file": "<path or name>",
        "status": "started"|"done"|"failed",
        "percent": float between 0-100,
        "elapsed": seconds for the file (float),
        "chunks": int (if done),
        "error": "<error str>" (if failed)
      }

    Returns:
      combined_chunks: list of chunk strings (all files)
      logs: dict of per-file logs & summary with timings
    """
    logs = {}
    combined_chunks = []

    def _worker(path):
        start = time.time()
        fname = os.path.basename(path)
        log = {"file": fname, "path": path, "start": start}

        try:
            chunks = extract_chunks(path)
            elapsed = time.time() - start
            log.update({"status": "done", "elapsed": elapsed, "chunks": len(chunks)})
            return path, chunks, log
        except Exception as e:
            elapsed = time.time() - start
            log.update({"status": "failed", "elapsed": elapsed, "error": str(e)})
            return path, [], log


    with ThreadPoolExecutor(max_workers=max_workers) as ex:
        futures = {ex.submit(_worker, p): p for p in file_paths}

        for fut in as_completed(futures):
            path, chunks, log = fut.result()

            # SAFE callback
            if progress_callback:
                progress_callback({
                    "file": os.path.basename(path),
                    "status": log.get("status"),
                    "chunks": log.get("chunks", 0),
                    "elapsed": log.get("elapsed", 0),
                    "error": log.get("error")
                })

            combined_chunks.extend(chunks)
            logs[os.path.basename(path)] = log

        
    logs["_summary"] = {
    "files": len(file_paths),
    "total_chunks": len(combined_chunks),
    "elapsed_total": sum(log["elapsed"] for log in logs.values() if "elapsed" in log)
    }
    
    return combined_chunks, logs 



# VECTOR DB (Embeddings + FAISS)

def build_vector_db(chunks, embed_model="sentence-transformers/all-MiniLM-L6-v2"):
    """
    Accepts list[str] chunks and builds FAISS (in-memory).
    Returns the FAISS vectorstore and timing/log info.
    """
    t0 = time.time()
    embedder = HuggingFaceEmbeddings(model_name=embed_model)

    # wrap as LangChain Documents expected by FAISS
    docs = [LCDocument(page_content=c) for c in chunks]

    t1 = time.time()
    vs = FAISS.from_documents(docs, embedder)
    t2 = time.time()

    logs = {
        "embedding_init_sec": t1 - t0,
        "faiss_build_sec": t2 - t1,
        "total_time_sec": t2 - t0,
        "num_documents": len(docs)
    }
    return vs, logs


# RAG CHAIN BUILDER

def build_rag_chain(retriever, api_key, model="llama-3.1-8b-instant", temperature=0.1):
    """
    Build and return a retrieval chain (combine_documents chain).
    """
    prompt = ChatPromptTemplate.from_template("""
Use ONLY the context below.

Context:
{context}

Instruction:
{input}

Response:
""")

    llm = ChatGroq(
        api_key=api_key,
        model=model,
        temperature=temperature
    )

    combine = create_stuff_documents_chain(llm=llm, prompt=prompt)

    return create_retrieval_chain(
        retriever=retriever,
        combine_docs_chain=combine
    )


# Small utility to write benchmark logs to disk

def save_benchmark_logs(logs, out_path="docex_benchmarks.json"):
    try:
        with open(out_path, "w", encoding="utf-8") as f:
            json.dump(logs, f, indent=2)
        return out_path
    except Exception as e:
        logger.error(f"Failed to save benchmark logs: {e}")
        return None

