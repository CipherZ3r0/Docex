from tracemalloc import start
import streamlit as st
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import io
from concurrent.futures import ThreadPoolExecutor
import time

# Import cleaning + chunking functions from text_processing.py
from text_processing import clean_text, chunk_text

# --- Set Tesseract path (adjust if needed) ---
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# --- Page config ---
st.set_page_config(page_title="Docex: Study Assistant", layout="wide")
st.title("📚 Docex — Your Study Assistant")

# ---- Text Extraction Functions ----
def extract_text_from_pdf(file):
    text = ""
    file_bytes = file.read()
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for i, page in enumerate(doc, start=1):
            page_text = page.get_text("text")
            if page_text.strip():
                text += f"\n--- Page {i} ---\n{page_text}"
            else:
                pix = page.get_pixmap(dpi=200)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                ocr_text = pytesseract.image_to_string(img)
                text += f"\n--- Page {i} (OCR) ---\n{ocr_text}"
    return text.strip()

def extract_text_from_docx(file):
    doc = Document(file)
    return "\n".join([p.text for p in doc.paragraphs]).strip()

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        if slide_text:
            text += f"\n--- Slide {i} ---\n" + "\n".join(slide_text)
    return text.strip()

def extract_text_from_txt(file):
    return file.read().decode("utf-8").strip()

def extract_text_from_image(file):
    img = Image.open(file)
    return pytesseract.image_to_string(img).strip()

# ---- Text Processing Function ----
def process_text(text):
    """
    Heavy text processing: cleaning + chunking
    """
    cleaned = clean_text(text)
    chunks = chunk_text(cleaned)
    return cleaned, chunks

# ---- Tabs ----
tab1, tab2 = st.tabs(["📂 Upload Materials", "❓ Ask Questions"])

# ----- Tab 1: Upload Materials -----
with tab1:
    st.header("Upload Your Study Materials")
    uploaded_files = st.file_uploader(
        "📤 Upload PDFs, DOCX, PPTX, TXT, or images:",
        type=["pdf", "docx", "pptx", "txt", "png", "jpg", "jpeg"],
        accept_multiple_files=True
    )
    uploaded_text = st.text_area("Or paste your text here:", height=200)

    all_texts = []

    start_time = time.perf_counter()

    # --- Extract from uploaded files ---
    if uploaded_files:
        st.success(f"✅ {len(uploaded_files)} file(s) uploaded successfully!")
        for file in uploaded_files:
            st.subheader(f"📂 {file.name}")
            try:
                extracted_text = ""
                if file.type == "application/pdf":
                    file.seek(0)
                    extracted_text = extract_text_from_pdf(file)
                elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    file.seek(0)
                    extracted_text = extract_text_from_docx(file)
                elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    file.seek(0)
                    extracted_text = extract_text_from_pptx(file)
                elif file.type == "text/plain":
                    file.seek(0)
                    extracted_text = extract_text_from_txt(file)
                elif file.type.startswith("image/"):
                    file.seek(0)
                    extracted_text = extract_text_from_image(file)
                else:
                    st.warning(f"⚠️ Unsupported file type: {file.type}")
                    continue

                if extracted_text:
                    all_texts.append(f"--- {file.name} ---\n{extracted_text}")
                    with st.expander(f"📝 View Extracted Text from {file.name}", expanded=False):
                        st.text_area(f"Extracted Text - {file.name}", extracted_text, height=300)
                else:
                    st.warning(f"⚠️ No text could be extracted from {file.name}.")

            except Exception as e:
                st.error(f"❌ Error extracting text from {file.name}: {e}")

    # --- Include pasted text ---
    if uploaded_text.strip():
        all_texts.append(f"--- Pasted Text ---\n{uploaded_text.strip()}")

    # --- Combine all texts ---
    combined_text = "\n\n".join(all_texts).strip() if all_texts else ""

    end_time = time.perf_counter()
    elapsed_time = end_time - start_time    
    st.write(f"⏱️ Time taken for extraction: {elapsed_time:.2f} seconds")

    print(f"⏱️ Text extraction took: {elapsed_time:.4f} seconds")
    st.info(f"⏱️ Text extraction took {elapsed_time:.2f} seconds")

    if combined_text:
        st.divider()
        with st.expander("📚 View Combined Study Material", expanded=True):
            st.subheader("📖 Combined Study Material")
            st.text_area("All Extracted & Pasted Text", combined_text, height=400)

    # --- Cleaning & Chunking ---
    if combined_text:
        st.divider()
        st.header("🧹 Text Cleaning & Chunking")
        if st.button("🧼 create memory"):
            with st.spinner("Cleaning and chunking text in background..."):
                with ThreadPoolExecutor() as executor:
                    future = executor.submit(process_text, combined_text)
                    cleaned, chunks = future.result()

            st.success(f"✅ Cleaned text and created {len(chunks)} chunks.")

            # with st.expander("🧾 View Cleaned Text", expanded=False):
            #     st.text_area("Cleaned Text", cleaned, height=300)

            # with st.expander("📦 View Text Chunks", expanded=False):
            #     for i, chunk in enumerate(chunks, start=1):
            #         st.markdown(f"**Chunk {i}**")
            #         st.text_area(f"Chunk {i}", chunk, height=200)

# ----- Tab 2: Ask Questions -----
with tab2:
    st.header("Ask Questions from Your Materials")
    question_input = st.text_area("Type your question here:")

    uploaded_question_file = st.file_uploader(
        "Or upload a question file (PDF, TXT, Image):",
        type=["pdf", "txt", "png", "jpg", "jpeg"]
    )

    if st.button("Get Answer"):
        st.info("🧠 Processing your question...")
        answer = "This is where your answer will appear once processing is implemented."
        st.success(answer)
