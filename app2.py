import time
from chunk_embedding import create_embeddings_from_text
import streamlit as st
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import io
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor, as_completed

# Import cleaning + chunking functions from text_processing.py
from text_processing import clean_text, chunk_text

# --- Set Tesseract path (adjust if needed) ---
# pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# --- Page config ---
st.set_page_config(page_title="Docex: Study Assistant", layout="wide")
st.title("Docex — Your Study Assistant")

# ---- Text Extraction Functions ----
def extract_text_from_pdf(file):
    text = ""
    file_bytes = file.read()
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for i, page in enumerate(doc, start=1):
            page_text = page.get_text("text")
            if page_text.strip():
                text += f"\n--- Page {i} ---\n{page_text}"
            else:
                pix = page.get_pixmap(dpi=250)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                ocr_text = pytesseract.image_to_string(img)
                text += f"\n--- Page {i} (OCR) ---\n{ocr_text}"
    return text.strip()

def extract_text_from_docx(file):
    doc = Document(file)
    return "\n".join([p.text for p in doc.paragraphs]).strip()

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        if slide_text:
            text += f"\n--- Slide {i} ---\n" + "\n".join(slide_text)
    return text.strip()

def extract_text_from_txt(file):
    return file.read().decode("utf-8").strip()

def extract_text_from_image(file):
    img = Image.open(file)
    return pytesseract.image_to_string(img).strip()

# ---- Text Processing Function ----
@st.cache_data(show_spinner=False)
def embedding_text(text):
    return create_embeddings_from_text(text)

# ---- Tabs ----
tab1, tab2 = st.tabs(["Upload Materials", "ask Questions"])

# ----- Tab 1: Upload Materials -----
with tab1:
    st.header("Upload Your Study Materials")
    uploaded_files = st.file_uploader(
        "Upload PDFs, DOCX, PPTX, TXT, or images:",
        type=["pdf", "docx", "pptx", "txt", "png", "jpg", "jpeg"],
        accept_multiple_files=True
    )
    uploaded_text = st.text_area("Or paste your text here:", height=200)

    all_texts = []

    # ---- Parallel File Extraction ----
    def process_file(file):
        try:
            if file.type == "application/pdf":
                file.seek(0)
                file_bytes = file.read()
                with fitz.open(stream=file_bytes, filetype="pdf") as doc:
                    text = ""
                    for i, page in enumerate(doc, start=1):
                        page_text = page.get_text("text")
                        if page_text.strip():
                            text += f"\n--- Page {i} ---\n{page_text}"
                        elif page.get_images():  # 🧠 only OCR if page has images
                            pix = page.get_pixmap(dpi=200)
                            img = Image.open(io.BytesIO(pix.tobytes("png")))
                            ocr_text = pytesseract.image_to_string(img)
                            text += f"\n--- Page {i} (OCR) ---\n{ocr_text}"
                    return file.name, text.strip()


            elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                file.seek(0)
                return file.name, extract_text_from_docx(file)

            elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                file.seek(0)
                return file.name, extract_text_from_pptx(file)

            elif file.type == "text/plain":
                file.seek(0)
                return file.name, extract_text_from_txt(file)

            elif file.type.startswith("image/"):
                file.seek(0)
                return file.name, extract_text_from_image(file)

            else:
                return file.name, None

        except Exception as e:
            return file.name, f"ERROR: {e}"
        
        
    
    start_time = time.perf_counter()
    if uploaded_files:
        st.success(f"✅ {len(uploaded_files)} file(s) uploaded successfully!")

        progress = st.empty()  # 🔹 feedback during extraction
        all_texts = []

        # 🔥 Run extraction in parallel (threads are fine for I/O)
        with ThreadPoolExecutor() as executor:
            results = list(executor.map(process_file, uploaded_files))

        progress.empty()

        for name, extracted_text in results:
            if extracted_text and not extracted_text.startswith("ERROR"):
                all_texts.append(f"--- {name} ---\n{extracted_text}")
                with st.expander(f"View Extracted Text from {name}", expanded=False):
                    st.text_area(f"Extracted Text - {name}", extracted_text, height=300)
            elif extracted_text and extracted_text.startswith("ERROR"):
                st.error(f"Error extracting text from {name}: {extracted_text}")
            else:
                st.warning(f"No text could be extracted from {name}.")


        # --- Include pasted text ---
        if uploaded_text.strip():
            all_texts.append(f"--- Pasted Text ---\n{uploaded_text.strip()}")

        # --- Combine all texts ---
        combined_text = "\n\n".join(all_texts).strip() if all_texts else ""

        if combined_text:
            st.divider()
            with st.expander("View Combined Study Material", expanded=True):
                st.subheader("Combined Study Material")
                st.text_area("All Extracted & Pasted Text", combined_text, height=400)

        end_time = time.perf_counter()
        elapsed_time = end_time - start_time
        print(f"Text extraction took: {elapsed_time:.4f} seconds")
        st.write(f"Time taken for extraction: {elapsed_time:.2f} seconds")

        # --- Cleaning & Chunking ---
        if combined_text:
            st.divider()
            st.header("Text Cleaning & Chunking")
            if st.button("create Brain"):
                with st.spinner("Cleaning and chunking text in background..."):
                   embedded_chunks = embedding_text(combined_text)

                st.success(f"Created {len(embedded_chunks)} embeddings successfully!")

                
                # with st.expander("🧾 View Cleaned Text", expanded=False):
                #     st.text_area("Cleaned Text", cleaned, height=300)

                # with st.expander("📦 View Text Chunks", expanded=False):
                #     for i, chunk in enumerate(chunks, start=1):
                #         st.markdown(f"**Chunk {i}**")
                #         st.text_area(f"Chunk {i}", chunk, height=200)

# ----- Tab 2: Ask Questions -----
with tab2:
    st.header("Ask Questions from Your Materials")
    question_input = st.text_area("Type your question here:")

    uploaded_question_file = st.file_uploader(
        "Or upload a question file (PDF, TXT, Image):",
        type=["pdf", "txt", "png", "jpg", "jpeg"]
    )

    if st.button("Get Answer"):
        st.info("Processing your question...")
        answer = "This is where your answer will appear once processing is implemented."
        st.success(answer)
